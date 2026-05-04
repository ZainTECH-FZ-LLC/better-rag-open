"""
Generate minimal placeholder template files for each department.

Run once during project setup:
    python src/templates/_generate_placeholders.py

Each file is a valid but empty-content Office document that the document
generators will use as a base template, filling in content from the LLM output.
Replace these with production-quality branded templates when available.
"""

from __future__ import annotations

import json
from pathlib import Path

HERE = Path(__file__).parent


def _load_style(dept: str) -> dict:
    style_path = HERE / dept / "_style.json"
    if style_path.exists():
        return json.loads(style_path.read_text())
    return {}


def _make_docx(path: Path, title: str, dept_display: str, style: dict) -> None:
    """Create a minimal branded .docx placeholder."""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        print(f"  [SKIP] python-docx not installed — {path.name}")
        return

    doc = Document()

    # Page margins
    docx_cfg = style.get("docx", {})
    for section in doc.sections:
        section.top_margin = Cm(docx_cfg.get("margin_top_cm", 2.54))
        section.bottom_margin = Cm(docx_cfg.get("margin_bottom_cm", 2.54))
        section.left_margin = Cm(docx_cfg.get("margin_left_cm", 3.18))
        section.right_margin = Cm(docx_cfg.get("margin_right_cm", 3.18))

    # Title
    heading = doc.add_heading(title, level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Subtitle
    sub = doc.add_paragraph(f"{dept_display} | BetterRAG Enterprise")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("")
    doc.add_paragraph(
        "This is a placeholder template. Replace with a production-quality "
        "branded document. The document generators will insert content above "
        "and below this placeholder text."
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    print(f"  [OK] {path.relative_to(HERE.parent.parent)}")


def _make_pptx(path: Path, title: str, dept_display: str, style: dict) -> None:
    """Create a minimal branded .pptx placeholder."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        print(f"  [SKIP] python-pptx not installed — {path.name}")
        return

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide

    primary_hex = style.get("colors", {}).get("primary", "#2E4057").lstrip("#")
    try:
        rgb = RGBColor(
            int(primary_hex[0:2], 16),
            int(primary_hex[2:4], 16),
            int(primary_hex[4:6], 16),
        )
    except (ValueError, IndexError):
        rgb = RGBColor(0x2E, 0x40, 0x57)

    if slide.shapes.title:
        slide.shapes.title.text = title
        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = rgb

    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = (
            f"{dept_display} | BetterRAG Enterprise\n"
            "Placeholder — replace with production template"
        )

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    print(f"  [OK] {path.relative_to(HERE.parent.parent)}")


def _make_xlsx(path: Path, title: str, dept_display: str, style: dict) -> None:
    """Create a minimal branded .xlsx placeholder."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        print(f"  [SKIP] openpyxl not installed — {path.name}")
        return

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]

    xlsx_cfg = style.get("xlsx", {})
    header_fill_hex = xlsx_cfg.get("header_fill", "#2E4057").lstrip("#")

    try:
        fill = PatternFill(
            start_color=header_fill_hex,
            end_color=header_fill_hex,
            fill_type="solid",
        )
    except Exception:
        fill = PatternFill(start_color="2E4057", end_color="2E4057", fill_type="solid")

    header_font = Font(
        name="Calibri",
        bold=True,
        color=xlsx_cfg.get("header_font_color", "FFFFFF").lstrip("#"),
        size=12,
    )

    ws["A1"] = title
    ws["A1"].font = header_font
    ws["A1"].fill = fill
    ws["A1"].alignment = Alignment(horizontal="center")

    ws["A2"] = dept_display
    ws["A3"] = "BetterRAG Enterprise — Placeholder Template"

    ws.column_dimensions["A"].width = 40

    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))
    print(f"  [OK] {path.relative_to(HERE.parent.parent)}")


TEMPLATES: list[tuple[str, str, str, str]] = [
    # (dept_dir, filename, title, display_name)
    # shared
    ("shared", "generic_report.docx",        "Generic Report",         "General"),
    ("shared", "generic_presentation.pptx",  "Generic Presentation",   "General"),
    ("shared", "generic_spreadsheet.xlsx",   "Generic Spreadsheet",    "General"),
    # hr
    ("hr",     "policy_report.docx",         "HR Policy Report",       "Human Resources"),
    ("hr",     "onboarding_deck.pptx",        "Onboarding Deck",        "Human Resources"),
    ("hr",     "headcount_tracker.xlsx",      "Headcount Tracker",      "Human Resources"),
    # finance
    ("finance","financial_report.docx",       "Financial Report",       "Finance"),
    ("finance","quarterly_review.pptx",       "Quarterly Review",       "Finance"),
    ("finance","budget_template.xlsx",        "Budget Template",        "Finance"),
    # sales
    ("sales",  "proposal.docx",              "Sales Proposal",         "Sales"),
    ("sales",  "sales_deck.pptx",            "Sales Deck",             "Sales"),
    ("sales",  "pipeline_tracker.xlsx",       "Pipeline Tracker",       "Sales"),
    # marketing
    ("marketing","campaign_brief.docx",      "Campaign Brief",         "Marketing"),
    ("marketing","brand_deck.pptx",          "Brand Deck",             "Marketing"),
    ("marketing","campaign_metrics.xlsx",    "Campaign Metrics",       "Marketing"),
]


def main() -> None:
    print("Generating placeholder templates in src/templates/\n")

    for dept, filename, title, display in TEMPLATES:
        style = _load_style(dept)
        path = HERE / dept / filename
        ext = path.suffix.lower()

        if path.exists():
            print(f"  [EXISTS] {path.relative_to(HERE.parent.parent)}")
            continue

        if ext == ".docx":
            _make_docx(path, title, display, style)
        elif ext == ".pptx":
            _make_pptx(path, title, display, style)
        elif ext == ".xlsx":
            _make_xlsx(path, title, display, style)

    print("\nDone. Replace placeholder files with production-quality branded templates.")


if __name__ == "__main__":
    main()
