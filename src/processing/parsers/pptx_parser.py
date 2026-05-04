"""PPTX parser — python-pptx for slide structure, speaker notes, shapes, images, and charts."""

from __future__ import annotations

import asyncio
import io

import structlog

from src.processing.parsers.base import DocumentParser, ParsedDocument

logger = structlog.get_logger()


class PPTXParser(DocumentParser):
    async def parse(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        return await asyncio.to_thread(self._parse_sync, file_bytes, filename)

    def _parse_sync(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(io.BytesIO(file_bytes))

        slides = []
        text_parts = []
        total_images = 0
        total_charts = 0

        for i, slide in enumerate(prs.slides):
            slide_data = {
                "index": i + 1,
                "title": "",
                "content": [],
                "notes": "",
                "shapes": [],
                "images": [],   # raw image bytes for vision extraction
                "charts": [],   # chart data extracted from chart objects
            }

            for shape in slide.shapes:
                shape_info = {"type": shape.shape_type, "name": shape.name}

                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_data["content"].append(text)
                        shape_info["text"] = text

                if hasattr(shape, "title") and shape == slide.shapes.title:
                    if shape.has_text_frame:
                        slide_data["title"] = shape.text_frame.text.strip()

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    shape_info["table"] = rows
                    table_text = "\n".join(
                        " | ".join(cell for cell in row if cell) for row in rows if any(row)
                    )
                    if table_text:
                        slide_data["content"].append(table_text)

                # Extract embedded images (charts/graphs pasted as pictures)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_blob = shape.image.blob
                        if img_blob and len(img_blob) > 5120:  # skip tiny icons (<5KB)
                            slide_data["images"].append(img_blob)
                            total_images += 1
                    except Exception:
                        pass

                # Extract chart data from native PowerPoint chart objects
                if shape.has_chart:
                    chart_data = _extract_chart_data(shape.chart)
                    if chart_data:
                        slide_data["charts"].append(chart_data)
                        total_charts += 1
                        # Also add chart as text content for non-vision fallback
                        chart_text = _chart_to_text(chart_data)
                        if chart_text:
                            slide_data["content"].append(chart_text)

                slide_data["shapes"].append(shape_info)

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                slide_data["notes"] = notes

            slides.append(slide_data)

            # Build text representation for the slide
            slide_text = f"--- Slide {i + 1} ---\n"
            if slide_data["title"]:
                slide_text += f"Title: {slide_data['title']}\n"
            for content in slide_data["content"]:
                slide_text += f"{content}\n"
            if slide_data["notes"]:
                slide_text += f"Speaker Notes: {slide_data['notes']}\n"
            text_parts.append(slide_text)

        full_text = "\n\n".join(text_parts)

        logger.info(
            "pptx_parser.parsed",
            filename=filename,
            slides=len(slides),
            images=total_images,
            charts=total_charts,
        )

        return ParsedDocument(
            text=full_text,
            slides=slides,
            page_count=len(slides),
            word_count=len(full_text.split()),
        )


def _extract_chart_data(chart) -> dict | None:
    """Extract data series from a native PowerPoint chart object."""
    try:
        from lxml import etree

        chart_data = {
            "type": str(chart.chart_type) if chart.chart_type else "unknown",
            "title": "",
            "categories": [],
            "series": [],
        }

        # Chart title
        if chart.has_title and chart.chart_title and chart.chart_title.has_text_frame:
            chart_data["title"] = chart.chart_title.text_frame.text.strip()

        # Categories (x-axis labels)
        try:
            plot = chart.plots[0]
            cats = plot.categories
            if cats:
                chart_data["categories"] = [str(c) for c in cats]
        except Exception:
            pass

        # Data series — extract from XML since .tx isn't available on all series types
        nsmap = {
            "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        chart_xml = chart._chartSpace.xml if hasattr(chart, "_chartSpace") else None

        for idx, series in enumerate(chart.series):
            # Get series name from XML (the .tx element)
            s_name = f"Series {idx + 1}"
            if chart_xml:
                try:
                    root = etree.fromstring(chart_xml) if isinstance(chart_xml, (str, bytes)) else chart_xml
                    # Find all ser elements and get the tx/v text
                    ser_elements = root.findall(".//c:ser", nsmap)
                    if idx < len(ser_elements):
                        tx_v = ser_elements[idx].find(".//c:tx//c:v", nsmap)
                        if tx_v is not None and tx_v.text:
                            s_name = tx_v.text
                except Exception:
                    pass

            s_data = {"name": s_name, "values": []}
            try:
                s_data["values"] = [float(v) if v is not None else None for v in series.values]
            except Exception:
                pass
            if s_data["values"]:
                chart_data["series"].append(s_data)

        return chart_data if chart_data["series"] else None
    except Exception as e:
        logger.debug("pptx_parser.chart_extraction_failed", error=str(e))
        return None


def _chart_to_text(chart_data: dict) -> str:
    """Render chart data as readable text."""
    parts = []

    chart_type = chart_data.get("type", "chart")
    title = chart_data.get("title", "")
    if title:
        parts.append(f"[Chart: {chart_type}] {title}")
    else:
        parts.append(f"[Chart: {chart_type}]")

    categories = chart_data.get("categories", [])
    series_list = chart_data.get("series", [])

    if not series_list:
        return "\n".join(parts)

    # Use categories if available, otherwise generate index labels
    num_points = max(len(s.get("values", [])) for s in series_list)
    if not categories:
        categories = [str(i + 1) for i in range(num_points)]

    # Build a markdown table
    header = "| # | " + " | ".join(s["name"] or f"Series {i+1}" for i, s in enumerate(series_list)) + " |"
    separator = "|" + "|".join("---" for _ in range(len(series_list) + 1)) + "|"
    parts.append(header)
    parts.append(separator)

    for j in range(min(len(categories), num_points)):
        row_vals = []
        for s in series_list:
            vals = s.get("values", [])
            val = vals[j] if j < len(vals) and vals[j] is not None else ""
            row_vals.append(str(val))
        parts.append(f"| {categories[j]} | " + " | ".join(row_vals) + " |")

    return "\n".join(parts)
