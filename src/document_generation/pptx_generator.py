"""PPTX generator — creates PowerPoint presentations from structured specs."""

from __future__ import annotations

import asyncio
import uuid
from pathlib import Path

import structlog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from src.document_generation.base import DocumentGenerator, DocumentSpec, GeneratedDocument

logger = structlog.get_logger()


class PPTXGenerator(DocumentGenerator):
    """
    Generates PowerPoint presentations using python-pptx.

    Slide types:
    - title: Title slide with subtitle
    - content: Title + bullet points
    - two_column: Two-column layout
    - chart: Title + embedded chart (via matplotlib)
    - table: Title + data table
    - section: Section divider slide
    """

    async def generate(self, spec: DocumentSpec, output_dir: Path) -> GeneratedDocument:
        return await asyncio.to_thread(self._generate_sync, spec, output_dir)

    def _generate_sync(self, spec: DocumentSpec, output_dir: Path) -> GeneratedDocument:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(prs, spec)

        # Content slides
        for section in spec.sections:
            slide_type = section.get("type", "content")

            if slide_type == "content":
                self._add_content_slide(prs, section)
            elif slide_type == "two_column":
                self._add_two_column_slide(prs, section)
            elif slide_type == "table":
                self._add_table_slide(prs, section)
            elif slide_type == "section":
                self._add_section_slide(prs, section)
            elif slide_type == "chart":
                self._add_chart_slide(prs, section, output_dir)
            else:
                self._add_content_slide(prs, section)

        # Save
        filename = f"{_sanitize(spec.title)}_{uuid.uuid4().hex[:8]}.pptx"
        filepath = output_dir / filename
        prs.save(str(filepath))

        size = filepath.stat().st_size
        logger.info("pptx_generator.saved", filename=filename, size=size, slides=len(prs.slides))

        return GeneratedDocument(
            filename=filename,
            filepath=filepath,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            file_type="pptx",
            size_bytes=size,
            metadata={"slide_count": len(prs.slides)},
        )

    def _add_title_slide(self, prs: Presentation, spec: DocumentSpec) -> None:
        layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        if title:
            title.text = spec.title
            for para in title.text_frame.paragraphs:
                para.font.size = Pt(36)
                para.font.bold = True

        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = spec.data.get("subtitle", "")

    def _add_content_slide(self, prs: Presentation, section: dict) -> None:
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        if title:
            title.text = section.get("title", "")

        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            bullets = section.get("bullets", [])
            for i, bullet in enumerate(bullets):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()

                para.text = bullet
                para.font.size = Pt(18)
                para.level = 0

                # Sub-bullets
                sub_bullets = section.get("sub_bullets", {}).get(str(i), [])
                for sub in sub_bullets:
                    sub_para = tf.add_paragraph()
                    sub_para.text = sub
                    sub_para.font.size = Pt(14)
                    sub_para.level = 1

    def _add_two_column_slide(self, prs: Presentation, section: dict) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        if title:
            title.text = section.get("title", "")

        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(5.8), Inches(5)
        )
        ltf = left_box.text_frame
        ltf.word_wrap = True
        for bullet in section.get("left_column", []):
            para = ltf.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(16)

        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(7), Inches(1.8), Inches(5.8), Inches(5)
        )
        rtf = right_box.text_frame
        rtf.word_wrap = True
        for bullet in section.get("right_column", []):
            para = rtf.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(16)

    def _add_table_slide(self, prs: Presentation, section: dict) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        if title:
            title.text = section.get("title", "")

        headers = section.get("headers", [])
        rows_data = section.get("rows", [])

        if not headers or not rows_data:
            return

        row_count = len(rows_data) + 1  # +1 for header
        col_count = len(headers)

        table_shape = slide.shapes.add_table(
            row_count, col_count,
            Inches(0.5), Inches(2), Inches(12), Inches(0.4 * row_count + 0.5)
        )
        table = table_shape.table

        # Header row
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = str(header)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)
                para.font.bold = True

        # Data rows
        for i, row in enumerate(rows_data):
            for j, value in enumerate(row):
                if j < col_count:
                    cell = table.cell(i + 1, j)
                    cell.text = str(value)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(11)

    def _add_section_slide(self, prs: Presentation, section: dict) -> None:
        layout = prs.slide_layouts[2]  # Section Header
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        if title:
            title.text = section.get("title", "")

    def _add_chart_slide(self, prs: Presentation, section: dict, output_dir: Path) -> None:
        """Add a slide with a chart rendered via matplotlib and inserted as an image."""
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        if title:
            title.text = section.get("title", "")

        chart_spec = section.get("chart", {})
        chart_type = chart_spec.get("type", "bar")
        labels = chart_spec.get("labels", [])
        values = chart_spec.get("values", [])

        if not labels or not values:
            return

        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt

            fig, ax = plt.subplots(figsize=(10, 5))

            if chart_type == "bar":
                ax.bar(labels, values, color="#4472C4")
            elif chart_type == "line":
                ax.plot(labels, values, marker="o", color="#4472C4", linewidth=2)
            elif chart_type == "pie":
                ax.pie(values, labels=labels, autopct="%1.1f%%")

            ax.set_title(section.get("title", ""), fontsize=14)
            plt.tight_layout()

            chart_path = output_dir / f"chart_{uuid.uuid4().hex[:8]}.png"
            fig.savefig(str(chart_path), dpi=150, bbox_inches="tight")
            plt.close(fig)

            slide.shapes.add_picture(
                str(chart_path), Inches(1.5), Inches(2), Inches(10), Inches(5)
            )
            chart_path.unlink(missing_ok=True)

        except Exception as e:
            logger.warn("pptx_generator.chart_failed", error=str(e))


def _sanitize(text: str) -> str:
    """Sanitize text for use as filename."""
    return "".join(c if c.isalnum() or c in " -_" else "" for c in text).strip()[:50]
