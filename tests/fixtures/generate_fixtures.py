"""Generate minimal valid test fixture files for all supported document types."""

from __future__ import annotations

from pathlib import Path

FIXTURES_DIR = Path(__file__).parent


def generate_pdf() -> None:
    """Generate a minimal 2-page PDF with text content."""
    try:
        from fpdf import FPDF
    except ImportError:
        print("fpdf2 not installed — skipping PDF fixture (pip install fpdf2)")
        return

    pdf = FPDF()
    pdf.set_margins(20, 20, 20)

    # Page 1
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 20)
    pdf.cell(0, 10, "BetterRAG Test Document", ln=True)
    pdf.ln(5)
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 8, "Section 1: HR Policy Overview", ln=True)
    pdf.set_font("Helvetica", size=11)
    pdf.multi_cell(
        0, 7,
        "This document describes the employee leave policy effective January 2024.\n\n"
        "Full-time employees are entitled to 20 days of paid time off (PTO) per year. "
        "Part-time employees receive PTO on a pro-rated basis.\n\n"
        "Parental leave: 16 weeks paid for primary caregivers, 4 weeks for secondary caregivers.",
    )
    pdf.ln(5)

    # Simple table
    pdf.set_font("Helvetica", "B", 11)
    col_widths = [60, 50, 60]
    headers = ["Leave Type", "Duration", "Eligibility"]
    for w, h in zip(col_widths, headers):
        pdf.cell(w, 8, h, border=1)
    pdf.ln()
    pdf.set_font("Helvetica", size=10)
    rows = [
        ("Annual PTO", "20 days/year", "All FT employees"),
        ("Sick Leave", "10 days/year", "All employees"),
        ("Parental Leave", "16 weeks", "Primary caregiver"),
    ]
    for row in rows:
        for w, v in zip(col_widths, row):
            pdf.cell(w, 7, v, border=1)
        pdf.ln()

    # Page 2
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 8, "Section 2: Financial Summary", ln=True)
    pdf.set_font("Helvetica", size=11)
    pdf.multi_cell(
        0, 7,
        "Q3 2024 Revenue: $14.1M (+22% YoY)\n"
        "EBITDA Margin: 31%\n"
        "Headcount: 245 (+18 vs Q2)\n\n"
        "The company exceeded revenue targets for the third consecutive quarter. "
        "EMEA contributed 37% of total revenue.",
    )

    pdf.output(str(FIXTURES_DIR / "sample.pdf"))
    print("Generated sample.pdf")


def generate_pptx() -> None:
    """Generate a minimal 3-slide PPTX."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("python-pptx not installed — skipping PPTX fixture (pip install python-pptx)")
        return

    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Q3 2024 Performance Review"
    slide.placeholders[1].text = "BetterRAG Enterprise — October 2024"

    # Slide 2: Bullets
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Highlights"
    tf = slide.placeholders[1].text_frame
    tf.text = "Revenue: $14.1M (+22% YoY)"
    for bullet in ["EBITDA margin: 31%", "NPS score: 68", "Headcount: 245"]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # Slide 3: Table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Revenue by Region"
    tbl = slide.shapes.add_table(4, 3, Inches(1), Inches(1.5), Inches(8), Inches(3)).table
    headers = ["Region", "Revenue", "vs Target"]
    for col, hdr in enumerate(headers):
        tbl.cell(0, col).text = hdr
    rows = [("Americas", "$6.1M", "+11%"), ("EMEA", "$5.2M", "+4%"), ("APAC", "$2.8M", "-7%")]
    for row_idx, (region, rev, vs) in enumerate(rows, 1):
        tbl.cell(row_idx, 0).text = region
        tbl.cell(row_idx, 1).text = rev
        tbl.cell(row_idx, 2).text = vs

    prs.save(str(FIXTURES_DIR / "sample.pptx"))
    print("Generated sample.pptx")


def generate_docx() -> None:
    """Generate a minimal DOCX with headings and a table."""
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        print("python-docx not installed — skipping DOCX fixture (pip install python-docx)")
        return

    doc = Document()
    doc.add_heading("BetterRAG Enterprise — HR Policy Document", 0)
    doc.add_heading("Section 1: Leave Policy", 1)
    doc.add_paragraph(
        "Full-time employees receive 20 days of paid time off per year. "
        "Leave balances reset on January 1st."
    )
    doc.add_heading("1.1 Parental Leave", 2)
    doc.add_paragraph(
        "Primary caregivers: 16 weeks fully paid.\n"
        "Secondary caregivers: 4 weeks fully paid."
    )

    # Table
    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    hdr_cells = table.rows[0].cells
    for i, h in enumerate(["Leave Type", "Duration", "Paid"]):
        hdr_cells[i].text = h
    rows = [
        ("Annual PTO", "20 days", "Yes"),
        ("Sick Leave", "10 days", "Yes"),
        ("Parental", "16 weeks", "Yes"),
    ]
    for row_idx, (lt, dur, paid) in enumerate(rows, 1):
        table.rows[row_idx].cells[0].text = lt
        table.rows[row_idx].cells[1].text = dur
        table.rows[row_idx].cells[2].text = paid

    doc.add_heading("Section 2: Remote Work", 1)
    doc.add_paragraph("Employees may work remotely up to 3 days per week.")

    doc.save(str(FIXTURES_DIR / "sample.docx"))
    print("Generated sample.docx")


def generate_xlsx() -> None:
    """Generate a minimal XLSX with two sheets."""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill
    except ImportError:
        print("openpyxl not installed — skipping XLSX fixture (pip install openpyxl)")
        return

    wb = openpyxl.Workbook()

    # Sheet 1: Revenue data
    ws1 = wb.active
    ws1.title = "Revenue"
    headers = ["Quarter", "Region", "Revenue", "Target", "Variance"]
    for col, h in enumerate(headers, 1):
        cell = ws1.cell(row=1, column=col, value=h)
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="1E4D8C")
        cell.font = Font(bold=True, color="FFFFFF")

    data = [
        ("Q3 2024", "Americas", 6100000, 5500000),
        ("Q3 2024", "EMEA", 5200000, 5000000),
        ("Q3 2024", "APAC", 2800000, 3000000),
    ]
    for row_idx, (q, region, rev, target) in enumerate(data, 2):
        ws1.cell(row=row_idx, column=1, value=q)
        ws1.cell(row=row_idx, column=2, value=region)
        ws1.cell(row=row_idx, column=3, value=rev)
        ws1.cell(row=row_idx, column=4, value=target)
        ws1.cell(row=row_idx, column=5, value=f"=C{row_idx}-D{row_idx}")

    # Sheet 2: Summary formulas
    ws2 = wb.create_sheet("Summary")
    ws2["A1"] = "Metric"
    ws2["B1"] = "Value"
    ws2["A2"] = "Total Revenue"
    ws2["B2"] = "=SUM(Revenue!C:C)"
    ws2["A3"] = "Total Target"
    ws2["B3"] = "=SUM(Revenue!D:D)"

    wb.save(str(FIXTURES_DIR / "sample.xlsx"))
    print("Generated sample.xlsx")


if __name__ == "__main__":
    print(f"Generating fixtures in: {FIXTURES_DIR}")
    generate_pdf()
    generate_pptx()
    generate_docx()
    generate_xlsx()
    print("Done.")
